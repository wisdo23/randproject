#!/usr/bin/env python3
import json, os, textwrap
from pptx import Presentation
from pptx.util import Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

BASE_DIR = os.path.dirname(__file__)
SLIDES_JSON = os.path.join(BASE_DIR, "slides_content.json")
OUT_PPTX = os.path.join(BASE_DIR, "presentation.pptx")
OUT_PDF = os.path.join(BASE_DIR, "presentation.pdf")

def load_slides(path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def make_pptx(slides, out_path):
    prs = Presentation()
    for slide in slides:
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get("title","")
        tx = s.shapes.placeholders[1].text_frame
        tx.clear()
        for b in slide.get("bullets",[]):
            p = tx.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)
    prs.save(out_path)
    print("Wrote PPTX ->", out_path)

def make_pdf(slides, out_path):
    c = canvas.Canvas(out_path, pagesize=letter)
    width, height = letter
    margin = 72
    for slide in slides:
        title = slide.get("title","")
        bullets = slide.get("bullets",[])
        c.setFont("Helvetica-Bold", 20)
        c.drawString(margin, height - margin - 10, title)
        c.setFont("Helvetica", 12)
        y = height - margin - 40
        wrapper = textwrap.TextWrapper(width=100)
        for b in bullets:
            lines = wrapper.wrap(b)
            for ln in lines:
                c.drawString(margin + 10, y, u"• " + ln)
                y -= 18
                if y < margin + 40:
                    c.showPage()
                    y = height - margin - 40
                    c.setFont("Helvetica", 12)
        c.showPage()
    c.save()
    print("Wrote PDF ->", out_path)

def main():
    slides = load_slides(SLIDES_JSON)
    make_pptx(slides, OUT_PPTX)
    make_pdf(slides, OUT_PDF)

if __name__ == "__main__":
    main()
